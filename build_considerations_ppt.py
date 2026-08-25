"""Build the editable 'AI Reporting — Key Considerations' slide (HOC template style)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
BLUE   = RGBColor(0x3E, 0x6D, 0xA8)
MAROON = RGBColor(0x7B, 0x2D, 0x3A)
BROWN  = RGBColor(0x8B, 0x5A, 0x1A)
GOLD   = RGBColor(0xE0, 0xA9, 0x3B)
CARD   = RGBColor(0xF4, 0xF6, 0xFA)
DARK   = RGBColor(0x2A, 0x2A, 0x2A)
GREY   = RGBColor(0x5A, 0x6A, 0x75)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD9, 0xE0, 0xEA)
TXBLUE = RGBColor(0x1F, 0x4E, 0x8C)
TXRED  = RGBColor(0xC8, 0x10, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def textbox(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    return tf


def run(p, text, size, color, bold=False, italic=False, name="Calibri"):
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    r.font.bold = bold; r.font.italic = italic; r.font.name = name
    return r


def rect(shape, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


# ---- header ----
tf = textbox(0.4, 0.26, 11.6, 0.55)
p = tf.paragraphs[0]
run(p, "HOUSE OF COLOUR  |  ", 22, NAVY, bold=True)
run(p, "AI Reporting — Key Considerations", 22, NAVY, bold=True)

tf = textbox(12.1, 0.26, 0.85, 0.5)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run(p, "T", 22, TXBLUE, bold=True, italic=True)
run(p, "x", 22, TXRED, bold=True, italic=True)

# gold rule
rect(MSO_SHAPE.RECTANGLE, 0.4, 0.92, 12.53, 0.05, GOLD)

# section title
tf = textbox(0.4, 1.05, 12.5, 0.55)
run(tf.paragraphs[0], "What to consider as AI usage grows", 26, NAVY, bold=True)

# intro
tf = textbox(0.4, 1.62, 12.5, 0.5)
run(tf.paragraphs[0],
    "The more we rely on AI for reporting, the more we expose our environment to external LLM APIs. "
    "These four areas need to be thought through — alongside the build, not after.",
    12, GREY)

# ---- pillars ----
pillars = [
    ("Security", NAVY, "DATA EXPOSURE", [
        "What's exposed during each API call?",
        "What extra is needed to protect our PII data?",
        "Where does data go — residency & retention?",
        "Who can access which data?",
        "Prompt injection via web mode",
    ]),
    ("Cost", BLUE, "USAGE & SPEND", [
        "Tokens cost money — schema sent each call",
        "How do we cap / limit AI usage?",
        "Higher cost per report",
        "Cost changes when the model changes",
        "Budget & quota planning",
    ]),
    ("Accuracy", MAROON, "RELIABILITY", [
        "How do we catch hallucinations / wrong numbers?",
        "How will testing & validation be done?",
        "How will debugging be carried out for AI reports?",
        "Output / experience shifts if model changes",
        "How do we keep metadata in sync?",
    ]),
    ("Governance", BROWN, "AUDIT & CONTROL", [
        "Audit trail for every AI query",
        "Which prompt & model version ran it?",
        "Can we reproduce a past result?",
        "Who owns the number?",
        "Compliance (GDPR / UK-DPA)",
    ]),
]

card_y, card_h, card_w, gap, x0 = 2.25, 4.55, 2.945, 0.252, 0.4
band_h = 0.5

for i, (name, color, band, items) in enumerate(pillars):
    x = x0 + i * (card_w + gap)
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h, CARD, line=LINE)
    rect(MSO_SHAPE.RECTANGLE, x, card_y, card_w, 0.08, color)            # top bar

    tf = textbox(x + 0.18, card_y + 0.16, card_w - 0.36, 0.45)
    run(tf.paragraphs[0], name, 17, color, bold=True)

    tf = textbox(x + 0.18, card_y + 0.70, card_w - 0.36, 0.3)
    run(tf.paragraphs[0], "WHAT TO CONSIDER", 9, color, bold=True)

    tf = textbox(x + 0.16, card_y + 1.04, card_w - 0.28, card_h - 1.04 - band_h - 0.06)
    for j, it in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        run(p, "•  ", 12, color, bold=True)
        run(p, it, 11.5, DARK)

    band_y = card_y + card_h - band_h
    bb = rect(MSO_SHAPE.RECTANGLE, x, band_y, card_w, band_h, color)
    btf = bb.text_frame; btf.word_wrap = True; btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    run(bp, band, 13, WHITE, bold=True)

# footer
tf = textbox(0.4, 7.0, 12.5, 0.35)
run(tf.paragraphs[0],
    "* These considerations scale with AI usage and should be planned alongside the build, not after.",
    11, GREY, italic=True)

import os
out = os.path.join(os.path.dirname(__file__), "..", "..", "..", "HOC (House of Colors)",
                   "HOC_AI_Key_Considerations.pptx")
out = os.path.abspath(out)
prs.save(out)
print("saved:", out)
