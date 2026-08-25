"""Build the editable 'AI Reporting — Cost & Consumption' slide (HOC template style).

Top ~70%  : factors that affect the cost (the main content, full width).
Lower ~30%: indicative cost only (accurate up-front estimation isn't possible).
Bottom    : AI usage can be controlled.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x1F, 0x3A, 0x5F)
BLUE   = RGBColor(0x3E, 0x6D, 0xA8)
MAROON = RGBColor(0x7B, 0x2D, 0x3A)
GREEN  = RGBColor(0x2E, 0x6B, 0x4F)
GOLD   = RGBColor(0xE0, 0xA9, 0x3B)
GOLDTN = RGBColor(0xFB, 0xF1, 0xD9)
BROWN  = RGBColor(0x8B, 0x5A, 0x1A)
CARD   = RGBColor(0xF4, 0xF6, 0xFA)
DARK   = RGBColor(0x2A, 0x2A, 0x2A)
GREY   = RGBColor(0x5A, 0x6A, 0x75)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xD9, 0xE0, 0xEA)
TINT_B = RGBColor(0xE8, 0xEE, 0xF6)
TINT_M = RGBColor(0xF6, 0xEC, 0xEF)
SEP    = RGBColor(0xCE, 0xD7, 0xE4)
TXBLUE = RGBColor(0x1F, 0x4E, 0x8C)
TXRED  = RGBColor(0xC8, 0x10, 0x2E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def tb(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    f = b.text_frame
    f.word_wrap = True; f.vertical_anchor = anchor
    f.margin_left = Inches(0.05); f.margin_right = Inches(0.05)
    f.margin_top = Inches(0.02); f.margin_bottom = Inches(0.02)
    return f


def run(p, t, s, c, bold=False, italic=False, name="Calibri"):
    r = p.add_run(); r.text = t
    r.font.size = Pt(s); r.font.color.rgb = c
    r.font.bold = bold; r.font.italic = italic; r.font.name = name
    return r


def rect(shape, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def cell(tbl, r, c, text, size, color, fill, bold=False, align=PP_ALIGN.LEFT):
    cl = tbl.cell(r, c)
    cl.fill.solid(); cl.fill.fore_color.rgb = fill
    cl.vertical_anchor = MSO_ANCHOR.MIDDLE
    cl.margin_left = Inches(0.08); cl.margin_right = Inches(0.05)
    cl.margin_top = Inches(0.01); cl.margin_bottom = Inches(0.01)
    cl.text = text
    p = cl.text_frame.paragraphs[0]; p.alignment = align
    rr = p.runs[0]; rr.font.size = Pt(size); rr.font.color.rgb = color
    rr.font.bold = bold; rr.font.name = "Calibri"


def column(x, y, w, h, color, title, desc, bullets):
    f = tb(x, y, w, h, MSO_ANCHOR.MIDDLE)
    p = f.paragraphs[0]; run(p, title, 12.5, color, bold=True)
    p = f.add_paragraph(); p.space_after = Pt(7); run(p, desc, 9, GREY, italic=True)
    for t in bullets:
        p = f.add_paragraph(); p.space_after = Pt(8)
        run(p, "•  ", 10.5, color, bold=True); run(p, t, 10.5, DARK)


# ---- header ----
f = tb(0.35, 0.18, 11.6, 0.5); p = f.paragraphs[0]
run(p, "HOUSE OF COLOUR  |  ", 20, NAVY, bold=True)
run(p, "AI Reporting — Cost & Consumption", 20, NAVY, bold=True)
f = tb(12.15, 0.18, 0.83, 0.46); p = f.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run(p, "T", 20, TXBLUE, bold=True, italic=True); run(p, "x", 20, TXRED, bold=True, italic=True)
rect(MSO_SHAPE.RECTANGLE, 0.35, 0.72, 12.63, 0.045, GOLD)

run(tb(0.35, 0.8, 12.6, 0.42).paragraphs[0], "What affects the cost of running the AI?", 23, NAVY, bold=True)
run(tb(0.35, 1.23, 12.6, 0.28).paragraphs[0],
    "AI-call cost is driven by the factors below. Any figures further down are indicative only — "
    "accurate up-front estimation is not possible.", 10.5, GREY)

# ============ TOP ~70%: factors (full width) ============
fx, fy, fw, fh = 0.35, 1.55, 12.63, 3.5
rect(MSO_SHAPE.ROUNDED_RECTANGLE, fx, fy, fw, fh, CARD, line=LINE)
rect(MSO_SHAPE.RECTANGLE, fx, fy, fw, 0.08, BLUE)
run(tb(fx + 0.2, fy + 0.13, fw - 0.4, 0.32).paragraphs[0],
    "Factors that drive the cost", 13.5, BLUE, bold=True)

cy0, cyh = fy + 0.55, fh - 0.66
cw = (fw - 0.4 - 3 * 0.0) / 4
xs = [fx + 0.2 + i * cw for i in range(4)]
# vertical separators
for i in range(1, 4):
    rect(MSO_SHAPE.RECTANGLE, xs[i] - 0.02, cy0 + 0.05, 0.012, cyh - 0.1, SEP)

column(xs[0] + 0.02, cy0, cw - 0.18, cyh, MAROON, "1 · Tokens per request",
       "what we send & get back, each call", [
           "Input: instructions + schema + history",
           "Output: SQL + chart + insights",
           "Several calls per question + retries",
           "Follow-ups carry past turns",
       ])
column(xs[1] + 0.02, cy0, cw - 0.18, cyh, NAVY, "2 · Model & settings",
       "which model, and how it's tuned", [
           "Model tier — Haiku < Sonnet < Opus",
           "Reasoning / “thinking” depth (billed)",
           "Max response length allowed",
       ])
column(xs[2] + 0.02, cy0, cw - 0.18, cyh, GREEN, "3 · Usage pattern",
       "who uses it, and how much", [
           "Number of users & concurrency",
           "Volume — questions / user / day",
           "User type & proficiency (trial-and-error)",
       ])
column(xs[3] + 0.02, cy0, cw - 0.18, cyh, BROWN, "4 · Query & data",
       "what's asked, over what schema", [
           "Query complexity & multi-step",
           "Schema size — NOT number of rows",
           "Research mode — web fee + bigger input",
       ])

# ============ LOWER ~30%: indicative cost ============
iy = fy + fh + 0.14
p = tb(0.35, iy, 12.63, 0.3).paragraphs[0]
run(p, "Estimated cost — INDICATIVE ONLY  ", 12.5, MAROON, bold=True)
run(p, "(accurate up-front estimation is not possible — it depends on the factors above)", 10, GREY, italic=True)

ty = iy + 0.34
# per-question table (left)
ql = [("Per question · Sonnet 4.6", "Tokens (in / out)", "≈ Cost", True),
      ("Typical (simple–medium)", "8–12K / ~1K", "$0.03–0.06", False),
      ("Complex / retry / research", "higher", "~$0.10–0.13", False)]
qt = slide.shapes.add_table(3, 3, Inches(0.35), Inches(ty), Inches(6.6), Inches(0.95)).table
qt.first_row = False; qt.horz_banding = False
qt.columns[0].width = Inches(2.9); qt.columns[1].width = Inches(2.1); qt.columns[2].width = Inches(1.6)
for ri, (a, b, c, hdr) in enumerate(ql):
    qt.rows[ri].height = Inches(0.31)
    if hdr:
        for ci, txt in enumerate((a, b, c)):
            cell(qt, ri, ci, txt, 9.8, WHITE, MAROON, bold=True,
                 align=PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
    else:
        cell(qt, ri, 0, a, 9.8, DARK, WHITE)
        cell(qt, ri, 1, b, 9.8, DARK, WHITE, align=PP_ALIGN.CENTER)
        cell(qt, ri, 2, c, 9.8, GREEN, TINT_M, bold=True, align=PP_ALIGN.CENTER)

# per-user/month table (right)
ux = 7.15
rr = [("Per user / day", "≈ per month", True),
      ("10", "~$11", False), ("25", "~$28", False), ("50", "~$55", False)]
ut = slide.shapes.add_table(4, 2, Inches(ux), Inches(ty), Inches(5.83), Inches(0.95)).table
ut.first_row = False; ut.horz_banding = False
ut.columns[0].width = Inches(3.2); ut.columns[1].width = Inches(2.63)
for ri, (a, b, hdr) in enumerate(rr):
    ut.rows[ri].height = Inches(0.225)
    if hdr:
        cell(ut, ri, 0, a, 9.5, WHITE, NAVY, bold=True)
        cell(ut, ri, 1, b, 9.5, WHITE, NAVY, bold=True, align=PP_ALIGN.CENTER)
    else:
        cell(ut, ri, 0, a, 9.5, DARK, WHITE)
        cell(ut, ri, 1, b, 9.5, NAVY, TINT_B, bold=True, align=PP_ALIGN.CENTER)
run(tb(ux, ty + 0.98, 5.83, 0.24).paragraphs[0],
    "~$0.05 / question · ~22 days · untrained users & research mode cost more.", 8, GREY, italic=True)

# ============ BOTTOM: controllable ============
by = 6.78
rect(MSO_SHAPE.ROUNDED_RECTANGLE, 0.35, by, 12.63, 0.42, GOLDTN, line=GOLD)
cf = tb(0.35, by + 0.02, 12.63, 0.38, MSO_ANCHOR.MIDDLE)
p = cf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run(p, "Most importantly — AI usage can be governed and controlled as it scales.", 12, BROWN, bold=True)

out = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..",
                                   "HOC (House of Colors)", "HOC_AI_Cost_Consumption.pptx"))
prs.save(out)
print("saved:", out)
