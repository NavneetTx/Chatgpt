"""Build the editable 'Demo Query List' slide (HOC template style)."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY = RGBColor(0x1F, 0x3A, 0x5F); BLUE = RGBColor(0x3E, 0x6D, 0xA8); MAROON = RGBColor(0x7B, 0x2D, 0x3A)
GOLD = RGBColor(0xE0, 0xA9, 0x3B); CARD = RGBColor(0xF4, 0xF6, 0xFA); DARK = RGBColor(0x2A, 0x2A, 0x2A)
GREY = RGBColor(0x5A, 0x6A, 0x75); WHITE = RGBColor(0xFF, 0xFF, 0xFF); LINE = RGBColor(0xD9, 0xE0, 0xEA)
TINT = {"NAVY": RGBColor(0xE9, 0xED, 0xF4), "BLUE": RGBColor(0xE8, 0xEE, 0xF6), "MAROON": RGBColor(0xF1, 0xE7, 0xEA)}
ASKTXT = RGBColor(0x22, 0x30, 0x4A)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])


def tb(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    b = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); f = b.text_frame
    f.word_wrap = True; f.vertical_anchor = anchor
    f.margin_left = Inches(0.06); f.margin_right = Inches(0.06); f.margin_top = Inches(0.03); f.margin_bottom = Inches(0.03)
    return f


def run(p, t, s, c, bold=False, italic=False, name="Calibri"):
    r = p.add_run(); r.text = t; r.font.size = Pt(s); r.font.color.rgb = c; r.font.bold = bold; r.font.italic = italic; r.font.name = name; return r


def rect(shape, x, y, w, h, color, line=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    if line is None: s.line.fill.background()
    else: s.line.color.rgb = line; s.line.width = Pt(0.75)
    s.shadow.inherit = False; return s


# header
f = tb(0.4, 0.26, 11.8, 0.55); p = f.paragraphs[0]
run(p, "HOUSE OF COLOUR  |  ", 22, NAVY, bold=True); run(p, "AI Reporting — Demo Query List", 22, NAVY, bold=True)
f = tb(12.1, 0.26, 0.85, 0.5); p = f.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
run(p, "T", 22, RGBColor(0x1F, 0x4E, 0x8C), bold=True, italic=True); run(p, "x", 22, RGBColor(0xC8, 0x10, 0x2E), bold=True, italic=True)
rect(MSO_SHAPE.RECTANGLE, 0.4, 0.92, 12.53, 0.05, GOLD)
run(tb(0.4, 1.05, 12.5, 0.5).paragraphs[0], "Sample queries we'll demonstrate", 25, NAVY, bold=True)
run(tb(0.4, 1.6, 12.5, 0.4).paragraphs[0],
    "Each demo opens with a chart, then keeps the conversation going — drill-downs and follow-ups, just like a chat.", 12, GREY)

queries = [
    ("Query 1", NAVY, "NAVY", "STACKED BAR",
     "Stacked bar of monthly revenue (Jan–Dec), each bar split into 4 territory segments, with the total on top.",
     ["Show only UK from this chart",
      "Which month had the highest revenue?",
      "Which territory grew the most vs last year?",
      "September revenue breakdown by territory",
      "Which territory is declining month over month?"]),
    ("Query 2", BLUE, "BLUE", "DONUT",
     "Donut chart — lead breakdown by source: Organic, Digital, B2B, Paid. Show count & % for each.",
     ["Show only Organic — which region did they come from?",
      "How many Digital leads converted to sessions?",
      "Which source has the highest conversion rate?"]),
    ("Query 3", MAROON, "MAROON", "LINE",
     "Quarterly bookings volume — trend line chart.",
     ["Break down Q1 by month",
      "Which session type drove the most bookings in Q2?",
      "Which quarter had the biggest growth vs last year?"]),
]

card_y, card_h, card_w, gap, x0 = 2.2, 4.55, 3.978, 0.3, 0.4
for i, (name, color, tintkey, ctype, ask, follows) in enumerate(queries):
    x = x0 + i * (card_w + gap)
    rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, card_y, card_w, card_h, CARD, line=LINE)
    rect(MSO_SHAPE.RECTANGLE, x, card_y, card_w, 0.08, color)
    # header: query name
    run(tb(x + 0.18, card_y + 0.16, 2.0, 0.4).paragraphs[0], name, 16, color, bold=True)
    # chart-type chip (right)
    chip_w = 1.55
    chip = rect(MSO_SHAPE.ROUNDED_RECTANGLE, x + card_w - chip_w - 0.18, card_y + 0.2, chip_w, 0.3, color)
    cf = chip.text_frame; cf.vertical_anchor = MSO_ANCHOR.MIDDLE; cp = cf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    run(cp, ctype, 9, WHITE, bold=True)
    # ask box
    abox = rect(MSO_SHAPE.ROUNDED_RECTANGLE, x + 0.18, card_y + 0.62, card_w - 0.36, 0.92, TINT[tintkey])
    af = abox.text_frame; af.word_wrap = True; af.vertical_anchor = MSO_ANCHOR.MIDDLE
    af.margin_left = Inches(0.1); af.margin_right = Inches(0.1)
    run(af.paragraphs[0], ask, 11, ASKTXT)
    # follow-ups label
    run(tb(x + 0.18, card_y + 1.62, card_w - 0.36, 0.3).paragraphs[0], "FOLLOW-UPS", 9, color, bold=True)
    # bullets
    bf = tb(x + 0.16, card_y + 1.95, card_w - 0.30, card_h - 1.95 - 0.12)
    for j, it in enumerate(follows):
        p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
        p.space_after = Pt(8)
        run(p, "›  ", 13, color, bold=True); run(p, it, 11.5, DARK)

run(tb(0.4, 7.0, 12.5, 0.35).paragraphs[0],
    "* Each follow-up continues the same chat thread; every answer stays scoped to the signed-in user.", 11, GREY, italic=True)

out = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "..", "HOC (House of Colors)", "HOC_AI_Demo_Queries.pptx"))
prs.save(out); print("saved:", out)
